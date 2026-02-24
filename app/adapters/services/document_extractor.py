import os
from typing import Optional

import fitz  # PyMuPDF
import pandas as pd
import pytesseract
from docx import Document
from PIL import Image
from pptx import Presentation
from transformers import AutoTokenizer

from app.application.interfaces import IDocumentExtractorInterface
from app.core.config import settings
from app.domain.entities.document_entity import ExtractedChunk


class DocumentExtractor(IDocumentExtractorInterface):
    def __init__(self, tokenizer_model: Optional[str] = None) -> None:
        model_name = tokenizer_model or getattr(settings, "OLLAMA_TOKENIZER_MODEL", "bert-base-uncased")
        self._tokenizer = AutoTokenizer.from_pretrained(model_name)
        # We only use the tokenizer for chunking, not inference.
        # Override model_max_length to suppress the spurious
        # "Token indices sequence length is longer than …" warning.
        self._tokenizer.model_max_length = int(1e30)

    async def extract(self, file_path: str) -> tuple[str, Optional[int]]:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self._extract_pdf(file_path)

        if ext in [".png", ".jpg", ".jpeg"]:
            return self._extract_image(file_path)

        if ext == ".csv":
            return self._extract_csv(file_path)

        if ext in [".xlsx", ".xls"]:
            return self._extract_excel(file_path)

        if ext == ".pptx":
            return self._extract_pptx(file_path)

        if ext == ".docx":
            return self._extract_docx(file_path)

        if ext == ".txt":
            return self._extract_txt(file_path)

        raise ValueError(f"Unsupported file type: {ext}")

    
    
    def chunk(self, text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
        """Split *text* into overlapping token-based chunks using the transformers tokeniser."""
        token_ids = self._tokenizer.encode(text, add_special_tokens=False)

        chunks: list[str] = []
        start = 0
        while start < len(token_ids):
            end = min(start + chunk_size, len(token_ids))
            chunk_text = self._tokenizer.decode(token_ids[start:end], skip_special_tokens=True)
            chunks.append(chunk_text)
            if end == len(token_ids):
                break
            start += chunk_size - overlap

        return chunks

    async def extract_and_chunk(
        self, file_path: str, chunk_size: int = 500, overlap: int = 50
    ) -> ExtractedChunk:
        """Extract text from *file_path* and return token-chunked results."""
        text, total_pages = await self.extract(file_path)
        chunks = self.chunk(text, chunk_size=chunk_size, overlap=overlap)
        return ExtractedChunk(chunks=chunks, total_pages=total_pages)


    # -----------------------------
    # PDF
    # -----------------------------

    def _extract_pdf(self, file_path: str) -> tuple[str, int]:
        doc = fitz.open(file_path)
        text = ""
        page_count = 0

        for page in doc:
            page_count += 1
            page_text = page.get_text()
            if not isinstance(page_text, str):
                page_text = str(page_text)

            # If PDF has no text (scanned)
            if not page_text.strip():
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                page_text = pytesseract.image_to_string(img)

            text += page_text + "\n"

        return text, page_count

    # -----------------------------
    # Images
    # -----------------------------

    def _extract_image(self, file_path: str) -> tuple[str, int]:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text, 1

    # -----------------------------
    # CSV
    # -----------------------------

    def _extract_csv(self, file_path: str) -> tuple[str, int]:
        df = pd.read_csv(file_path)
        return df.to_string(), 1

    # -----------------------------
    # Excel
    # -----------------------------

    def _extract_excel(self, file_path: str) -> tuple[str, int]:
        df = pd.read_excel(file_path, sheet_name=None)
        text = ""
        total_sheets = 0
        for sheet, data in df.items():
            text += f"\nSheet: {sheet}\n"
            text += data.to_string()
            total_sheets += 1
        return text, total_sheets

    # -----------------------------
    # PPTX
    # -----------------------------

    def _extract_pptx(self, file_path: str) -> tuple[str, int]:
        prs = Presentation(file_path)
        text = ""
        
        total_slides = len(prs.slides)
        for slide in prs.slides:
            for shape in slide.shapes:
                text_frame = getattr(shape, "text_frame", None)
                if text_frame is not None:
                    text += text_frame.text + "\n"
                    
        return text, total_slides
    
    # -----------------------------
    # DOCX
    # -----------------------------

    def _extract_docx(self, file_path: str) -> tuple[str, int]:
        doc = Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
        return text, 1
    
    
    # -----------------------------
    # TXT
    # -----------------------------

    def _extract_txt(self, file_path: str) -> tuple[str, int]:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        return text, 1

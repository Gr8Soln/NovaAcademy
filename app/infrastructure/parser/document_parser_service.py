"""Document parser — extracts text from PDF, DOCX, PPTX, images, TXT."""

from __future__ import annotations

import os
from typing import List

from app.domain.exceptions import InvalidFileTypeError
from app.interfaces.services.document_parser_service import IDocumentParserService, ParsedChunk


class LocalDocumentParserService(IDocumentParserService):
    """Parse documents using open-source libraries (pypdf, python-docx, python-pptx, pytesseract)."""

    CHUNK_SIZE = 1000  # characters per chunk
    CHUNK_OVERLAP = 200

    async def parse(self, file_path: str, file_type: str) -> List[ParsedChunk]:
        if file_type == "pdf":
            return await self._parse_pdf(file_path)
        elif file_type == "docx":
            return await self._parse_docx(file_path)
        elif file_type == "pptx":
            return await self._parse_pptx(file_path)
        elif file_type == "txt":
            return await self._parse_txt(file_path)
        elif file_type == "image":
            return await self._parse_image(file_path)
        else:
            raise InvalidFileTypeError(f"Unsupported file type: {file_type}")

    def _chunk_text(self, text: str, page_number: int | None = None) -> List[ParsedChunk]:
        """Split text into overlapping chunks."""
        chunks: List[ParsedChunk] = []
        start = 0
        while start < len(text):
            end = start + self.CHUNK_SIZE
            chunk_text = text[start:end]
            if chunk_text.strip():
                chunks.append(ParsedChunk(content=chunk_text.strip(), page_number=page_number))
            start += self.CHUNK_SIZE - self.CHUNK_OVERLAP
        return chunks

    async def _parse_pdf(self, file_path: str) -> List[ParsedChunk]:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        all_chunks: List[ParsedChunk] = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            all_chunks.extend(self._chunk_text(text, page_number=i))
        return all_chunks

    async def _parse_docx(self, file_path: str) -> List[ParsedChunk]:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return self._chunk_text(full_text)

    async def _parse_pptx(self, file_path: str) -> List[ParsedChunk]:
        from pptx import Presentation

        prs = Presentation(file_path)
        all_chunks: List[ParsedChunk] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            slide_text = "\n".join(texts)
            if slide_text.strip():
                all_chunks.extend(self._chunk_text(slide_text, page_number=i))
        return all_chunks

    async def _parse_txt(self, file_path: str) -> List[ParsedChunk]:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        return self._chunk_text(text)

    async def _parse_image(self, file_path: str) -> List[ParsedChunk]:
        try:
            import pytesseract
            from PIL import Image

            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            return self._chunk_text(text) if text.strip() else []
        except ImportError:
            return [ParsedChunk(content="[Image OCR not available — install pytesseract]")]
